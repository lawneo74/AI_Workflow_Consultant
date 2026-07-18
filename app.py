"""AI Workflow Architect — intelligent Prompt & Workflow Generator.

A Streamlit app that turns a raw task description into a chained, copy-paste
workflow. The user picks which AI tools they actually have access to
(Perplexity AI, Claude, ChatGPT, NotebookLM, Gemini incl. Nano Banana), and a
three-tier Anthropic backend does the rest:

1. Router  — Haiku classifies the task's complexity and flags whether live
   web research would improve the plan.
2. Research (optional) — when the router asks for it and a Perplexity API key
   is configured, one sonar-pro search grounds the plan in current facts.
3. Generator — Haiku (simple) or Sonnet (complex) designs the workflow,
   applying professional prompt-engineering principles to every step.
4. Reviewer — Sonnet does a final quality pass that checks the plan is
   aligned with the user's goal and refines the prompts and tool routing
   before anything is shown to the user.

Each step is labeled with the exact app to paste it into; steps routed to
Claude also carry a recommended Claude model and effort level. When a
workflow spans multiple apps a Transition note explains how to carry the
data across. The finished plan can be downloaded as Markdown, Word (.docx),
or PDF.

The app never executes the generated prompts — it is strictly a copy-paste
generator. (The optional Perplexity call gathers planning context before
generation; it never runs the generated prompts.)
"""

import hmac
import io
import json
import os

import anthropic
import requests
import streamlit as st

ROUTER_MODEL = "claude-haiku-4-5"
SIMPLE_GENERATOR_MODEL = "claude-haiku-4-5"
COMPLEX_GENERATOR_MODEL = "claude-opus-4-8"
REVIEW_MODEL = "claude-sonnet-5"
PERPLEXITY_RESEARCH_MODEL = "sonar-pro"
PERPLEXITY_API_URL = "https://api.perplexity.ai/chat/completions"

# Rich tool-capabilities knowledge base, maintained as Markdown in the repo
# (editable directly on GitHub — see the note at the top of the file). Loaded
# once per script run; if missing or unreadable the planner falls back to the
# built-in TOOL_CATALOG descriptions.
TOOL_CAPABILITIES_FILE = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "tool_capabilities.md"
)


def load_tool_capabilities() -> str | None:
    try:
        with open(TOOL_CAPABILITIES_FILE, encoding="utf-8") as f:
            text = f.read().strip()
        return text or None
    except OSError:
        return None


TOOL_CAPABILITIES = load_tool_capabilities()

# Desired output budget for the generator/reviewer. claude-sonnet-5 runs
# adaptive thinking by default, so those tokens share max_tokens with the JSON
# plan; a full multi-step plan with detailed prompts needs generous headroom or
# it truncates into invalid JSON. The SDK requires streaming for budgets this
# large (and streaming also avoids HTTP timeouts), so these calls stream and
# read the final message.
#
# This is a *ceiling to aim for*, not the value sent to the API: each model
# rejects (400) a max_tokens above its own output cap, so plan_budget() clamps
# it to MODEL_MAX_OUTPUT_TOKENS per model before every call.
PLAN_MAX_TOKENS = 160000

# Per-model maximum output tokens. The Messages API returns a 400
# invalid_request_error when max_tokens exceeds the model's cap.
MODEL_MAX_OUTPUT_TOKENS = {
    "claude-haiku-4-5": 64000,
    "claude-sonnet-5": 128000,
    "claude-opus-4-8": 128000,
}


def plan_budget(model: str) -> int:
    """max_tokens for a planner call, clamped to the model's output cap."""
    return min(PLAN_MAX_TOKENS, MODEL_MAX_OUTPUT_TOKENS.get(model, 64000))


def thinking_kwargs(model: str) -> dict:
    """Thinking config per model. claude-opus-4-8 runs without thinking unless
    adaptive thinking is requested explicitly; claude-sonnet-5 runs adaptive by
    default, and claude-haiku-4-5 uses the legacy budget style — so only Opus
    needs the explicit parameter."""
    if model == "claude-opus-4-8":
        return {"thinking": {"type": "adaptive"}}
    return {}


# --------------------------------------------------------------------------- #
# Context attachments (MD / DOCX / PPTX)                                       #
# --------------------------------------------------------------------------- #

ATTACHMENT_TYPES = ["md", "docx", "pptx"]
MAX_ATTACHMENT_CHARS = 20_000  # per file, after extraction
MAX_CONTEXT_CHARS = 60_000  # total across all attachments
ROUTER_CONTEXT_CHARS = 4_000  # the router only needs a taste for classification


def extract_attachment_text(name: str, data: bytes) -> str | None:
    """Best-effort plain-text extraction from an uploaded file.

    Returns None when the file can't be parsed — the caller warns and
    continues, matching the app's best-effort philosophy.
    """
    lower = name.lower()
    try:
        if lower.endswith(".md"):
            return data.decode("utf-8", errors="replace")
        if lower.endswith(".docx"):
            from docx import Document

            doc = Document(io.BytesIO(data))
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
            return "\n".join(parts)
        if lower.endswith(".pptx"):
            from pptx import Presentation

            prs = Presentation(io.BytesIO(data))
            parts = []
            for i, slide in enumerate(prs.slides, start=1):
                texts = [
                    shape.text_frame.text.strip()
                    for shape in slide.shapes
                    if shape.has_text_frame and shape.text_frame.text.strip()
                ]
                if slide.has_notes_slide:
                    note = slide.notes_slide.notes_text_frame.text.strip()
                    if note:
                        texts.append(f"(Speaker notes: {note})")
                if texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(parts)
    except ModuleNotFoundError:
        return None
    except Exception:
        return None
    return None


def build_attachment_context(files) -> tuple[str, list[str], list[str]]:
    """Extract text from uploaded files into one planning-context block.

    Returns (context_text, included_filenames, notes). Oversized content is
    truncated with a visible note; unreadable files warn and are skipped.
    """
    sections: list[str] = []
    included: list[str] = []
    notes: list[str] = []
    total = 0
    for f in files or []:
        text = extract_attachment_text(f.name, f.getvalue())
        if text is None or not text.strip():
            notes.append(f"⚠️ Couldn't read “{f.name}” — it was skipped.")
            continue
        text = text.strip()
        if len(text) > MAX_ATTACHMENT_CHARS:
            text = text[:MAX_ATTACHMENT_CHARS]
            notes.append(f"📎 “{f.name}” was long — only the first part was used.")
        if total + len(text) > MAX_CONTEXT_CHARS:
            text = text[: max(0, MAX_CONTEXT_CHARS - total)]
            if not text:
                notes.append(f"📎 “{f.name}” was skipped — context limit reached.")
                continue
            notes.append(f"📎 “{f.name}” was truncated — context limit reached.")
        total += len(text)
        included.append(f.name)
        sections.append(f'--- Attachment: "{f.name}" ---\n{text}')
    return "\n\n".join(sections), included, notes


def compose_task_context(task: str, attachment_context: str, limit: int | None = None) -> str:
    """The task plus attached reference material, as sent to the planner."""
    if not attachment_context:
        return task
    context = attachment_context if limit is None else attachment_context[:limit]
    return (
        f"{task}\n\n"
        "Reference material attached by the user (use it as authoritative "
        f"context for this task):\n\n{context}"
    )

# The catalog of tools the user may have. Each carries a routing description
# (fed to the model) and a badge color (for the UI / step labels).
TOOL_CATALOG = {
    "Perplexity AI": {
        "color": "#1F6F8B",
        "description": (
            "Real-time web search and current events. Best for live research, "
            "market/competitive scans, fact-checking, and answers backed by "
            "cited, up-to-date sources."
        ),
    },
    "Claude": {
        "color": "#C15F3C",
        "description": (
            "Long-form writing, nuanced analysis and synthesis, careful "
            "step-by-step reasoning, coding, and working with large documents. "
            "Best when quality of thinking and prose matters most."
        ),
    },
    "ChatGPT": {
        "color": "#10A37F",
        "description": (
            "Versatile general assistant with strong coding and debugging, data "
            "analysis, brainstorming, and image generation (DALL·E). A dependable "
            "all-rounder."
        ),
    },
    "NotebookLM": {
        "color": "#1A73E8",
        "description": (
            "Grounded question-answering over the user's own uploaded sources: "
            "faithful, source-cited summaries, study guides, and audio overviews. "
            "Best when there is a defined corpus of documents to reason over."
        ),
    },
    "Gemini (incl. Nano Banana)": {
        "color": "#9334E6",
        "description": (
            "Google-ecosystem tasks, strong multimodal understanding, and image "
            "generation/editing via Nano Banana. Best for visual content "
            "creation, image editing, and Google Workspace integration."
        ),
    },
}

ALL_TOOLS = list(TOOL_CATALOG)

ROUTER_SYSTEM = """\
You are a task-complexity router for an AI workflow planner. Classify the
user's task so the right generator model can be selected.

Classify as "complex" when the task involves any of: multi-step workflows,
software engineering or coding, research combined with synthesis, work that
spans more than one tool, or open-ended deliverables (reports, whitepapers,
applications). Classify as "simple" for single-shot tasks: a lookup, a quick
rewrite, a summary of provided text, a single well-scoped prompt.

Also decide whether a live web search would materially improve the workflow
plan — set "needs_research" true when the task depends on current facts,
recent tools/versions, market conditions, or anything likely to have changed
since your training data. When true, write a focused "research_query" (one
search query capturing what the planner needs to know). When false, use an
empty string for "research_query".
"""

ROUTER_SCHEMA = {
    "type": "object",
    "properties": {
        "complexity": {"type": "string", "enum": ["simple", "complex"]},
        "reasoning": {
            "type": "string",
            "description": "One-sentence justification for the classification.",
        },
        "needs_research": {
            "type": "boolean",
            "description": "True when live web research would improve the plan.",
        },
        "research_query": {
            "type": "string",
            "description": "Search query for the research step; empty if not needed.",
        },
    },
    "required": ["complexity", "reasoning", "needs_research", "research_query"],
    "additionalProperties": False,
}

# Professional prompt-engineering principles the generated prompts must embody.
PROMPT_ENGINEERING_PRINCIPLES = """\
You write every step "prompt" in the capacity of a senior professional prompt
engineer: each one must be a production-quality prompt the user can paste
verbatim and get an excellent result on the first try. Hold every prompt to
this bar:

- Role & context: open by assigning the target app a specific expert persona
  (not a generic "assistant") and give it the situational context it needs —
  who the work is for, why it matters, and what came before this step.
- Explicit deliverable: state the objective and the concrete deliverable
  precisely and unambiguously; never leave the app to guess scope.
- Inputs & references: name every input the step depends on — the user's
  attached reference material and artifacts produced by earlier steps — and
  say exactly how each should be used.
- Structure with delimiters: separate instructions from content with headings,
  numbered requirements, or delimiters (### sections, triple backticks) so the
  app can't confuse the two.
- Output specification: define the exact format, length, structure, and file
  type of the expected output, including section-by-section outlines for
  documents and column definitions for tabular work.
- Constraints & edge cases: list what to do, what to avoid, boundary
  conditions to respect, and assumptions to surface rather than silently make.
- Success criteria & self-check: end analytical prompts by stating how the
  result will be judged and asking the app to verify its output against those
  criteria before finishing.
- Reasoning guidance: for multi-step or analytical work, direct the app to
  plan or reason through the problem before producing the final answer.
- Tone & audience: specify the intended audience, register, and level of
  detail whenever the deliverable is user-facing prose.
Keep every prompt specific and self-contained — a reader with no other context
must be able to execute it. Never write a vague one-liner."""

# Rule applied by both the generator and the reviewer: Claude steps must carry
# a concrete model + effort recommendation.
CLAUDE_STEP_RULE = """\
- For every step routed to "Claude" you MUST recommend both the Claude model
  to use (in "model", e.g. claude-sonnet-5 for most work, claude-fable-5 or
  claude-opus-4-8 for the hardest reasoning) AND the "effort" to use ("Low",
  "Medium", or "High"). Effort is the reasoning depth the user should set:
  the /effort setting in Claude Code, or Extended Thinking in the Claude app
  when "High". Reserve "High" for the hardest reasoning, design, or debugging
  steps. For steps in other tools, "effort" must be an empty string (use
  "model" for a mode/model suggestion only when clearly useful)."""

# Rule applied by both the generator and the reviewer: transitions must give
# the user real, app-specific handoff guidance — and, for Claude Code steps,
# the one-time project setup they need before running the prompt.
TRANSITION_RULE = """\
- Transitions are the user's guidance for moving from the previous step into
  this one. Make them concrete and specific to the two apps involved — never a
  generic "paste the output". When the app changes, state exactly what to
  carry across and how: what to export or save, the file format to use, and
  how to hand it to the next app (paste it, upload/attach the file, add it as
  a source, or reference it). Fill "transition" whenever a handoff or setup
  needs explaining; use an empty string only when the app is unchanged and
  nothing needs carrying across.
- Claude Code setup: when a step uses Claude Code (app "Claude" with "model"
  set to "Claude Code"), its "transition" MUST first walk the user through the
  one-time project setup before they run the prompt — even on a first step,
  where the transition is the getting-started note:
    - create a project folder with a sensible structure for the task (name the
      key folders/files, e.g. src/, docs/, data/, tests/);
    - add a CLAUDE.md at the project root capturing the project context, goal,
      conventions, and constraints Claude Code should follow;
    - add any Skills or reference files that will help (e.g. drop the user's
      attached material into the project, or note a Skill to enable);
    - then open Claude Code in that folder and run the step's prompt.
  Keep this concise and actionable, tailored to the specific task."""


def _tool_menu(selected_tools: list[str]) -> str:
    """Render the description block for the tools the user selected."""
    return "\n".join(
        f'- "{name}": {TOOL_CATALOG[name]["description"]}' for name in selected_tools
    )


def _capabilities_block() -> str:
    """The maintained capabilities knowledge base, when available."""
    if not TOOL_CAPABILITIES:
        return ""
    return f"""
Use this maintained tool-capabilities knowledge base as your authoritative
reference for each tool's strengths, output formats, weaknesses, and proven
tool sequences. Note: finer-grained capabilities it mentions (Claude Code,
Claude Projects, Perplexity Deep Research / Computer / Spaces, Nano Banana)
are modes WITHIN the routable tools — a step's "app" must always be one of
the user's tools listed above, and the specific mode belongs in the step's
"model" field and prompt text (e.g. app "Claude" with model "Claude Code",
or app "Perplexity AI" with model "Deep Research").

<tool_capabilities>
{TOOL_CAPABILITIES}
</tool_capabilities>
"""


SIMPLICITY_RULE = """\
- Simplicity is the product's core promise: recommend the SIMPLEST workflow
  that fully delivers the goal. A single step in a single tool is the ideal
  outcome whenever it is genuinely sufficient — add steps or tools only when
  each one is necessary to reach the stated goal."""


def build_generator_system(selected_tools: list[str]) -> str:
    return f"""\
You are an AI workflow architect and senior professional prompt engineer. The
user has access to ONLY the following tools — you must never route a step to
any tool outside this list:

{_tool_menu(selected_tools)}
{_capabilities_block()}
Given a task, design the optimal chained workflow across these tools.

Rules:
{SIMPLICITY_RULE}
- Choose the fewest tools that genuinely fit the task. Use a multi-app chain
  only when different phases clearly belong in different tools.
- Every step must name the exact app to paste the prompt into (from the list
  above) and provide a complete, copy-paste-ready prompt for that app.
{CLAUDE_STEP_RULE}
{TRANSITION_RULE}
- "effort_level" reflects the user's expected hands-on effort: "Low",
  "Medium", or "High".
- If the user attached reference material, ground the plan in it and make the
  step prompts point the target app at the relevant attached content.
- If research findings are provided with the task, treat them as current,
  authoritative context: ground the strategy and step prompts in them
  (correct tool names, versions, and facts) instead of stale knowledge.
- Do not include account-signup steps, installation instructions, or
  usage-limit caveats, and never write a Perplexity MCP setup guide. (Claude
  Code *project* setup — CLAUDE.md, Skills, folder structure — is the one
  exception and IS wanted; it belongs in the transition per the rule above.)

{PROMPT_ENGINEERING_PRINCIPLES}"""


def build_candidates_system(selected_tools: list[str]) -> str:
    """System prompt for the complex path: multiple candidates, then select."""
    return f"""\
{build_generator_system(selected_tools)}

This task is non-trivial, so follow this two-phase procedure:

Phase 1 — Generate 2 to 3 GENUINELY DIFFERENT candidate workflows. Candidates
must differ in strategy or tool sequence (e.g. different lead tool, different
phase structure, single-tool vs. chained), not be cosmetic variants of one
plan. Every candidate must independently satisfy all the rules above and
fully deliver the user's goal.

Phase 2 — Select the winner: the SIMPLEST candidate that fully delivers the
goal. Goal completeness is a hard gate — never pick an incomplete plan just
because it is shorter. Among goal-complete candidates, prefer the fewest
steps, then the fewest tool switches, then the least hands-on user effort.
Set "selected_index" to the winner's zero-based position in "candidates" and
explain the choice in "selection_rationale" (1-2 sentences: why this one is
the simplest that gets the job done)."""


REVIEW_SYSTEM_TEMPLATE = """\
You are a senior professional prompt engineer performing the FINAL quality
pass on a drafted AI workflow before it reaches the user. The user has access
to ONLY these tools — every step must stay within this list:

{menu}
{capabilities}
Critically review the draft workflow and return an improved final version.
Check and fix:
- Goal alignment (most important): executing the workflow step by step must
  fully deliver the user's stated goal — nothing missing, no scope drift, no
  steps the user did not ask for, and the final step produces the intended
  deliverable. Rework the plan if it falls short.
- Simplicity: this must be the simplest workflow that fully delivers the
  goal. Remove any step or tool switch the goal does not require; merge steps
  that one tool can complete in a single prompt. Never simplify away
  something the goal needs.
- Tool routing: every step uses an allowed tool and the best-fit tool for its
  phase per the capabilities reference; the chain uses the fewest tools that
  genuinely fit.
- Claude recommendations: every step routed to "Claude" carries a concrete
  Claude model in "model" and an "effort" of "Low", "Medium", or "High" (the
  /effort setting in Claude Code, or Extended Thinking in the Claude app when
  "High"); steps in other tools keep "effort" as an empty string.
- Prompt craftsmanship: each prompt must meet the senior-prompt-engineer bar
  defined below — specific expert role, explicit deliverable, named inputs,
  delimited structure, exact output spec, constraints, and success criteria.
  Rewrite any prompt that falls short; a vague or generic prompt is a defect.
- Attachments: if the user attached reference material, the prompts must
  direct each app to the relevant attached content.
- Transitions: concrete and app-specific wherever context moves between
  apps — the user must know exactly what to carry across and how (export,
  save, paste, upload/attach, or add as a source). Every Claude Code step
  (app "Claude", model "Claude Code") must carry its one-time project setup
  in the transition: project folder structure, a CLAUDE.md capturing context
  and conventions, and any Skills or reference files. Empty only when nothing
  needs carrying across. Add or rewrite any transition that falls short.
- Coherence: steps flow logically and together fully accomplish the task.
- If research findings accompany the task, the plan must be consistent with
  them (current tool names, versions, and facts).

Rewrite and tighten prompts as needed. Then write a short "review_summary"
(2-3 sentences) that states explicitly whether the final plan is aligned with
the user's goal, confirms it is the simplest sufficient plan (mention when
alternative candidates were compared), and notes what you verified or
improved. Do not include account-signup steps, installation instructions, or
usage-limit caveats, and never write a Perplexity MCP setup guide — but keep
Claude Code *project* setup (CLAUDE.md, Skills, folder structure) in the
relevant transition.

{principles}"""


def _workflow_properties(selected_tools: list[str]) -> dict:
    """Shared JSON-schema properties for the generator and reviewer."""
    return {
        "strategy_summary": {
            "type": "string",
            "description": "Two or three sentences explaining the overall plan.",
        },
        "recommended_environments": {
            "type": "array",
            "items": {"type": "string", "enum": selected_tools},
        },
        "effort_level": {"type": "string", "enum": ["Low", "Medium", "High"]},
        "steps": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "title": {"type": "string"},
                    "app": {"type": "string", "enum": selected_tools},
                    "model": {
                        "type": "string",
                        "description": (
                            "Recommended mode or model within the app. Required "
                            "for Claude steps (a Claude model id); empty string "
                            "if not applicable."
                        ),
                    },
                    "effort": {
                        "type": "string",
                        "enum": ["Low", "Medium", "High", ""],
                        "description": (
                            "Reasoning effort for Claude steps (the /effort "
                            "setting in Claude Code; High = Extended Thinking "
                            "in the Claude app). Empty string for other tools."
                        ),
                    },
                    "transition": {
                        "type": "string",
                        "description": (
                            "Concrete, app-specific guidance for moving from "
                            "the previous step into this one: what to carry "
                            "across and how (export/save/paste/upload/attach). "
                            "For Claude Code steps, also the one-time project "
                            "setup (folder structure, CLAUDE.md, Skills). "
                            "Empty string when nothing needs carrying across."
                        ),
                    },
                    "prompt": {
                        "type": "string",
                        "description": "The copy-paste prompt for this step.",
                    },
                },
                "required": ["title", "app", "model", "effort", "transition", "prompt"],
                "additionalProperties": False,
            },
        },
    }


def build_generator_schema(selected_tools: list[str]) -> dict:
    return {
        "type": "object",
        "properties": _workflow_properties(selected_tools),
        "required": [
            "strategy_summary",
            "recommended_environments",
            "effort_level",
            "steps",
        ],
        "additionalProperties": False,
    }


def build_candidates_schema(selected_tools: list[str]) -> dict:
    """Schema for the complex path: 2-3 candidate workflows plus a selection."""
    return {
        "type": "object",
        "properties": {
            "candidates": {
                "type": "array",
                "items": build_generator_schema(selected_tools),
                "description": (
                    "2-3 genuinely different candidate workflows, each of "
                    "which fully delivers the user's goal."
                ),
            },
            "selected_index": {
                "type": "integer",
                "description": (
                    "Zero-based index of the simplest goal-complete candidate."
                ),
            },
            "selection_rationale": {
                "type": "string",
                "description": (
                    "1-2 sentences on why the selected candidate is the "
                    "simplest workflow that gets the job done."
                ),
            },
        },
        "required": ["candidates", "selected_index", "selection_rationale"],
        "additionalProperties": False,
    }


def build_review_schema(selected_tools: list[str]) -> dict:
    props = _workflow_properties(selected_tools)
    props["review_summary"] = {
        "type": "string",
        "description": (
            "2-3 sentences stating whether the plan is aligned with the goal "
            "and what was verified or improved."
        ),
    }
    return {
        "type": "object",
        "properties": props,
        "required": [
            "review_summary",
            "strategy_summary",
            "recommended_environments",
            "effort_level",
            "steps",
        ],
        "additionalProperties": False,
    }


def get_passcodes() -> list[str]:
    """Return the two configured login passcodes.

    Read from the APP_PASSCODES env var or Streamlit secrets as a
    comma-separated pair, e.g. APP_PASSCODES="alpha-1234,bravo-5678".
    """
    raw = os.environ.get("APP_PASSCODES")
    if not raw:
        try:
            raw = st.secrets.get("APP_PASSCODES", None)
        except Exception:
            raw = None
    if not raw:
        return []
    return [p.strip() for p in str(raw).split(",") if p.strip()][:2]


def check_login() -> bool:
    """Render the login page until a valid passcode is entered."""
    if st.session_state.get("authenticated"):
        return True

    st.title("🧭 AI Workflow Architect")
    st.subheader("Sign in")

    passcodes = get_passcodes()
    if not passcodes:
        st.warning("No passcodes configured. Set `APP_PASSCODES` to enable login.")
        return False

    with st.form("login_form"):
        entered = st.text_input("Passcode", type="password")
        submitted = st.form_submit_button("Sign in", type="primary", use_container_width=True)

    if submitted:
        if any(hmac.compare_digest(entered, code) for code in passcodes):
            st.session_state["authenticated"] = True
            st.rerun()
        else:
            st.error("Incorrect passcode.")
    return False


def get_client() -> anthropic.Anthropic | None:
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        try:
            api_key = st.secrets.get("ANTHROPIC_API_KEY", None)
        except Exception:
            api_key = None
    if not api_key:
        api_key = st.session_state.get("api_key_input") or None
    if not api_key:
        return None
    return anthropic.Anthropic(api_key=api_key)


def get_perplexity_key() -> str | None:
    api_key = os.environ.get("PERPLEXITY_API_KEY")
    if not api_key:
        try:
            api_key = st.secrets.get("PERPLEXITY_API_KEY", None)
        except Exception:
            api_key = None
    if not api_key:
        api_key = st.session_state.get("perplexity_key_input") or None
    return api_key


def run_research(api_key: str, query: str) -> str | None:
    """Run one Perplexity search to gather live context for the planner.

    Returns the research findings (with source URLs when available), or None
    if the search fails — research is best-effort and never blocks planning.
    """
    try:
        response = requests.post(
            PERPLEXITY_API_URL,
            headers={"Authorization": f"Bearer {api_key}"},
            json={
                "model": PERPLEXITY_RESEARCH_MODEL,
                "messages": [
                    {
                        "role": "system",
                        "content": (
                            "You are a research assistant. Answer concisely "
                            "with current, factual information the requester "
                            "can plan against. Include key names, versions, "
                            "and dates."
                        ),
                    },
                    {"role": "user", "content": query},
                ],
            },
            timeout=45,
        )
        response.raise_for_status()
        data = response.json()
        findings = data["choices"][0]["message"]["content"]
        sources = data.get("citations") or data.get("search_results") or []
        urls = [s["url"] if isinstance(s, dict) else s for s in sources][:5]
        if urls:
            findings += "\n\nSources:\n" + "\n".join(f"- {u}" for u in urls)
        return findings
    except (requests.RequestException, KeyError, IndexError, ValueError):
        return None


def extract_json(response) -> dict:
    text = next(b.text for b in response.content if b.type == "text")
    return json.loads(text)


def route_task(client: anthropic.Anthropic, task: str) -> dict:
    """Layer 1 — Haiku classifies task complexity and the need for research."""
    response = client.messages.create(
        model=ROUTER_MODEL,
        max_tokens=512,
        system=ROUTER_SYSTEM,
        output_config={"format": {"type": "json_schema", "schema": ROUTER_SCHEMA}},
        messages=[{"role": "user", "content": f"Task: {task}"}],
    )
    return extract_json(response)


def generate_workflow(
    client: anthropic.Anthropic,
    task: str,
    complexity: str,
    selected_tools: list[str],
    research: str | None = None,
) -> tuple[dict, str]:
    """Layer 2 — Haiku (simple) or Sonnet (complex) drafts the workflow."""
    model = (
        COMPLEX_GENERATOR_MODEL if complexity == "complex" else SIMPLE_GENERATOR_MODEL
    )
    content = (
        f"Available tools: {', '.join(selected_tools)}\n\n"
        f"Design the workflow for this task:\n\n{task}"
    )
    if research:
        content += f"\n\n---\n\nResearch findings (from a live web search):\n\n{research}"
    with client.messages.stream(
        model=model,
        max_tokens=plan_budget(model),
        system=build_generator_system(selected_tools),
        output_config={
            "format": {"type": "json_schema", "schema": build_generator_schema(selected_tools)}
        },
        messages=[{"role": "user", "content": content}],
        **thinking_kwargs(model),
    ) as stream:
        response = stream.get_final_message()
    return extract_json(response), model


def generate_candidate_workflows(
    client: anthropic.Anthropic,
    task: str,
    selected_tools: list[str],
    research: str | None = None,
) -> dict:
    """Layer 2 (complex path) — Opus drafts 2-3 candidates and selects the
    simplest one that fully delivers the goal."""
    model = COMPLEX_GENERATOR_MODEL
    content = (
        f"Available tools: {', '.join(selected_tools)}\n\n"
        f"Design candidate workflows for this task and select the simplest "
        f"one that gets the job done:\n\n{task}"
    )
    if research:
        content += f"\n\n---\n\nResearch findings (from a live web search):\n\n{research}"
    with client.messages.stream(
        model=model,
        max_tokens=plan_budget(model),
        system=build_candidates_system(selected_tools),
        output_config={
            "format": {"type": "json_schema", "schema": build_candidates_schema(selected_tools)}
        },
        messages=[{"role": "user", "content": content}],
        **thinking_kwargs(model),
    ) as stream:
        response = stream.get_final_message()
    return extract_json(response)


def pick_selected_candidate(payload: dict) -> tuple[dict, int, str]:
    """Validate the candidates payload and return (workflow, count, rationale).

    Raises KeyError/IndexError/TypeError on a malformed payload — callers
    treat that as "candidate comparison unavailable" and fall back.
    """
    candidates = payload["candidates"]
    if not isinstance(candidates, list) or not candidates:
        raise KeyError("candidates")
    index = payload.get("selected_index", 0)
    if not isinstance(index, int) or not (0 <= index < len(candidates)):
        index = 0
    workflow = candidates[index]
    if not isinstance(workflow, dict) or "steps" not in workflow:
        raise KeyError("steps")
    return workflow, len(candidates), str(payload.get("selection_rationale", ""))


def review_workflow(
    client: anthropic.Anthropic,
    task: str,
    selected_tools: list[str],
    draft: dict,
    research: str | None = None,
    selection_note: str | None = None,
) -> dict:
    """Layer 3 — Sonnet reviews goal alignment and returns the polished plan."""
    system = REVIEW_SYSTEM_TEMPLATE.format(
        menu=_tool_menu(selected_tools),
        capabilities=_capabilities_block(),
        principles=PROMPT_ENGINEERING_PRINCIPLES,
    )
    content = (
        f"Available tools: {', '.join(selected_tools)}\n\n"
        f"Original task:\n{task}\n\n"
        f"Draft workflow to review (JSON):\n{json.dumps(draft, indent=2)}"
    )
    if selection_note:
        content += (
            f"\n\n---\n\nThis draft was selected from multiple candidate "
            f"workflows as the simplest that delivers the goal. Selection "
            f"rationale: {selection_note}"
        )
    if research:
        content += f"\n\n---\n\nResearch findings (from a live web search):\n\n{research}"
    with client.messages.stream(
        model=REVIEW_MODEL,
        max_tokens=plan_budget(REVIEW_MODEL),
        system=system,
        output_config={
            "format": {"type": "json_schema", "schema": build_review_schema(selected_tools)}
        },
        messages=[{"role": "user", "content": content}],
    ) as stream:
        response = stream.get_final_message()
    return extract_json(response)


def app_badge(app_name: str) -> str:
    color = TOOL_CATALOG.get(app_name, {}).get("color", "#555555")
    return (
        f'<span style="background-color:{color}; color:white; padding:3px 12px; '
        f'border-radius:12px; font-size:0.85em; font-weight:600; '
        f'white-space:nowrap;">{app_name}</span>'
    )


def step_meta_text(step: dict) -> str:
    """Model/effort metadata suffix for a step, shared by UI and exports."""
    meta = ""
    if step.get("model", "").strip():
        meta += f"Suggested model/mode: {step['model']}"
    if step.get("effort", "").strip():
        meta += ("  ·  " if meta else "") + f"Effort: {step['effort']}"
    return meta


# --------------------------------------------------------------------------- #
# Downloadable exports                                                          #
# --------------------------------------------------------------------------- #

def build_markdown(task: str, workflow: dict, attachments: list[str] | None = None) -> str:
    lines = ["# AI Workflow Plan", "", f"**Task:** {task}", ""]
    if attachments:
        lines += [f"**Context files:** {', '.join(attachments)}", ""]
    lines += ["## Strategy", ""]
    lines.append(
        "**Recommended tools:** "
        + ", ".join(workflow.get("recommended_environments", []))
    )
    lines.append("")
    lines.append(f"**Effort level:** {workflow.get('effort_level', '')}")
    lines.append("")
    lines.append(workflow.get("strategy_summary", ""))
    lines.append("")

    if workflow.get("review_summary"):
        lines += ["## Final Review", "", workflow["review_summary"], ""]

    lines += ["## Chained Workflow", ""]
    for i, step in enumerate(workflow.get("steps", []), start=1):
        lines.append(f"### Step {i}: {step['title']}  —  [{step['app']}]")
        lines.append("")
        if step.get("transition", "").strip():
            lines.append(f"> **Transition:** {step['transition']}")
            lines.append("")
        meta = f"**Paste into:** {step['app']}"
        if step.get("model", "").strip():
            meta += f"  ·  **Suggested model/mode:** {step['model']}"
        if step.get("effort", "").strip():
            meta += f"  ·  **Effort:** {step['effort']}"
        lines.append(meta)
        lines.append("")
        lines.append("```text")
        lines.append(step["prompt"])
        lines.append("```")
        lines.append("")
    return "\n".join(lines)


def build_docx(task: str, workflow: dict, attachments: list[str] | None = None) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor

    doc = Document()
    doc.add_heading("AI Workflow Plan", level=0)

    p = doc.add_paragraph()
    p.add_run("Task: ").bold = True
    p.add_run(task)
    if attachments:
        p = doc.add_paragraph()
        p.add_run("Context files: ").bold = True
        p.add_run(", ".join(attachments))

    doc.add_heading("Strategy", level=1)
    p = doc.add_paragraph()
    p.add_run("Recommended tools: ").bold = True
    p.add_run(", ".join(workflow.get("recommended_environments", [])))
    p = doc.add_paragraph()
    p.add_run("Effort level: ").bold = True
    p.add_run(str(workflow.get("effort_level", "")))
    doc.add_paragraph(workflow.get("strategy_summary", ""))

    if workflow.get("review_summary"):
        doc.add_heading("Final Review", level=1)
        doc.add_paragraph(workflow["review_summary"])

    doc.add_heading("Chained Workflow", level=1)
    for i, step in enumerate(workflow.get("steps", []), start=1):
        doc.add_heading(f"Step {i}: {step['title']}  —  [{step['app']}]", level=2)
        if step.get("transition", "").strip():
            tp = doc.add_paragraph()
            tp.add_run("Transition: ").bold = True
            tp.add_run(step["transition"])
        meta = doc.add_paragraph()
        meta.add_run("Paste into: ").bold = True
        meta.add_run(step["app"])
        if step_meta_text(step):
            meta.add_run("   ·   ")
            meta.add_run(step_meta_text(step))
        # Prompt block in monospace.
        for line in step["prompt"].split("\n"):
            pp = doc.add_paragraph()
            run = pp.add_run(line if line else " ")
            run.font.name = "Courier New"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_pdf(task: str, workflow: dict, attachments: list[str] | None = None) -> bytes:
    from xml.sax.saxutils import escape

    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    styles = getSampleStyleSheet()
    code_style = ParagraphStyle(
        "PromptCode",
        parent=styles["Normal"],
        fontName="Courier",
        fontSize=8.5,
        leading=11,
        alignment=TA_LEFT,
        backColor="#F4F4F4",
        borderPadding=6,
        leftIndent=4,
        rightIndent=4,
    )

    def para(text: str, style_name: str = "Normal"):
        return Paragraph(escape(text).replace("\n", "<br/>"), styles[style_name])

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=LETTER,
        leftMargin=0.9 * inch, rightMargin=0.9 * inch,
        topMargin=0.9 * inch, bottomMargin=0.9 * inch,
    )
    flow = [para("AI Workflow Plan", "Title")]
    flow += [para(f"<b>Task:</b> {escape(task)}"), Spacer(1, 10)]
    if attachments:
        flow += [
            para(f"<b>Context files:</b> {', '.join(attachments)}"),
            Spacer(1, 10),
        ]

    flow.append(para("Strategy", "Heading1"))
    flow.append(
        Paragraph(
            "<b>Recommended tools:</b> "
            + escape(", ".join(workflow.get("recommended_environments", []))),
            styles["Normal"],
        )
    )
    flow.append(
        Paragraph(
            f"<b>Effort level:</b> {escape(str(workflow.get('effort_level', '')))}",
            styles["Normal"],
        )
    )
    flow.append(para(workflow.get("strategy_summary", "")))
    flow.append(Spacer(1, 8))

    if workflow.get("review_summary"):
        flow.append(para("Final Review", "Heading1"))
        flow.append(para(workflow["review_summary"]))
        flow.append(Spacer(1, 8))

    flow.append(para("Chained Workflow", "Heading1"))
    for i, step in enumerate(workflow.get("steps", []), start=1):
        flow.append(para(f"Step {i}: {step['title']}  —  [{step['app']}]", "Heading2"))
        if step.get("transition", "").strip():
            flow.append(
                Paragraph(
                    f"<b>Transition:</b> {escape(step['transition'])}", styles["Normal"]
                )
            )
        meta = f"<b>Paste into:</b> {escape(step['app'])}"
        if step_meta_text(step):
            meta += f"   ·   {escape(step_meta_text(step))}"
        flow.append(Paragraph(meta, styles["Normal"]))
        flow.append(Spacer(1, 4))
        flow.append(Paragraph(escape(step["prompt"]).replace("\n", "<br/>"), code_style))
        flow.append(Spacer(1, 12))

    doc.build(flow)
    return buf.getvalue()


def render_downloads(
    task: str, workflow: dict, attachments: list[str] | None = None
) -> None:
    st.markdown("**Download this plan**")
    col_md, col_docx, col_pdf = st.columns(3)
    with col_md:
        st.download_button(
            "⬇️ Markdown",
            data=build_markdown(task, workflow, attachments),
            file_name="ai_workflow_plan.md",
            mime="text/markdown",
            use_container_width=True,
        )
    with col_docx:
        try:
            docx_bytes = build_docx(task, workflow, attachments)
            st.download_button(
                "⬇️ Word (.docx)",
                data=docx_bytes,
                file_name="ai_workflow_plan.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True,
            )
        except ModuleNotFoundError:
            st.button("Word (.docx)", disabled=True, help="python-docx not installed", use_container_width=True)
    with col_pdf:
        try:
            pdf_bytes = build_pdf(task, workflow, attachments)
            st.download_button(
                "⬇️ PDF",
                data=pdf_bytes,
                file_name="ai_workflow_plan.pdf",
                mime="application/pdf",
                use_container_width=True,
            )
        except ModuleNotFoundError:
            st.button("PDF", disabled=True, help="reportlab not installed", use_container_width=True)


def render_workflow(
    task: str, workflow: dict, attachments: list[str] | None = None
) -> None:
    st.subheader("Strategy")
    col1, col2 = st.columns([3, 1])
    with col1:
        st.markdown(
            "**Recommended Tools:** "
            + " ".join(app_badge(env) for env in workflow["recommended_environments"]),
            unsafe_allow_html=True,
        )
    with col2:
        st.markdown(f"**Effort Level:** {workflow['effort_level']}")
    st.write(workflow["strategy_summary"])

    if workflow.get("review_summary"):
        st.success(f"**✅ Final review:** {workflow['review_summary']}")

    st.subheader("Chained Workflow")
    for i, step in enumerate(workflow["steps"], start=1):
        if step["transition"].strip():
            st.info(f"**Transition:** {step['transition']}")
        with st.container(border=True):
            header_col, badge_col = st.columns([3, 1])
            with header_col:
                st.markdown(f"#### Step {i}: {step['title']}")
            with badge_col:
                st.markdown(
                    f'<div style="text-align:right; padding-top:14px;">'
                    f'{app_badge(step["app"])}</div>',
                    unsafe_allow_html=True,
                )
            caption = f"Paste in: **{step['app']}**"
            if step.get("model", "").strip():
                caption += f" · Suggested model/mode: `{step['model']}`"
            if step.get("effort", "").strip():
                caption += f" · Effort: **{step['effort']}**"
            st.caption(caption)
            st.code(step["prompt"], language=None, wrap_lines=True)

    st.divider()
    render_downloads(task, workflow, attachments)


def start_new_session() -> None:
    """Clear the current plan and inputs, keeping auth and API key."""
    for key in (
        "workflow",
        "workflow_task",
        "workflow_notes",
        "workflow_attachments",
        "workflow_context_sig",
        "task_input",
    ):
        st.session_state.pop(key, None)


def main() -> None:
    st.set_page_config(page_title="AI Workflow Architect", page_icon="🧭", layout="centered")

    if not check_login():
        return

    with st.sidebar:
        st.markdown("### Session")
        if st.button("🔄 Start new session", use_container_width=True):
            start_new_session()
            st.rerun()
        if st.button("Sign out", use_container_width=True):
            st.session_state.clear()
            st.rerun()

    st.title("🧭 AI Workflow Architect")
    st.caption(
        "Describe your goal, pick the AI tools you have, and get a chained, "
        "copy-paste workflow with each step routed to the best app."
    )

    client = get_client()
    if client is None:
        st.text_input(
            "Anthropic API key",
            type="password",
            key="api_key_input",
            help="Used only to plan your workflow. Never stored.",
        )
        client = get_client()

    perplexity_key = get_perplexity_key()
    if perplexity_key is None:
        with st.expander("Live research (optional)"):
            st.text_input(
                "Perplexity API key",
                type="password",
                key="perplexity_key_input",
                help=(
                    "When set, the planner searches the web for current facts "
                    "before designing your workflow. Never stored."
                ),
            )
        perplexity_key = get_perplexity_key()

    selected_tools = st.multiselect(
        "Which AI tools do you have access to?",
        options=ALL_TOOLS,
        default=st.session_state.get("selected_tools", ["Claude", "Perplexity AI"]),
        key="selected_tools",
        help="Select at least one. Workflows are routed only to the tools you pick.",
    )

    task = st.text_area(
        "Describe your task or goal.",
        height=140,
        key="task_input",
        placeholder="e.g. Research AI agent frameworks in 2026 and write a whitepaper.",
    )

    uploaded_files = st.file_uploader(
        "Attach reference files (optional)",
        type=ATTACHMENT_TYPES,
        accept_multiple_files=True,
        key="context_files",
        help=(
            "Markdown, Word (.docx), or PowerPoint (.pptx) files whose "
            "content gives the planner context about your task."
        ),
    )
    context_sig = ";".join(
        f"{f.name}:{f.size}" for f in (uploaded_files or [])
    )

    # Editing the task or changing attachments invalidates the previous
    # result: clear it from memory so a stale workflow is never shown next to
    # different inputs.
    if "workflow" in st.session_state and (
        task.strip() != st.session_state.get("workflow_task")
        or context_sig != st.session_state.get("workflow_context_sig", "")
    ):
        del st.session_state["workflow"]
        for key in ("workflow_task", "workflow_notes", "workflow_attachments",
                    "workflow_context_sig"):
            st.session_state.pop(key, None)

    if st.button("Architect My Workflow", type="primary", use_container_width=True):
        if client is None:
            st.error("Enter your Anthropic API key above to continue.")
            st.stop()
        if not selected_tools:
            st.warning("Select at least one AI tool.")
            st.stop()
        if not task.strip():
            st.warning("Describe your task first.")
            st.stop()

        with st.spinner("Reading attached files…" if uploaded_files else "Preparing…"):
            attachment_context, attachment_names, notes = build_attachment_context(
                uploaded_files
            )
        if attachment_names:
            notes.append("📎 Planned with context from: " + ", ".join(attachment_names))
        planner_task = compose_task_context(task.strip(), attachment_context)
        router_task = compose_task_context(
            task.strip(), attachment_context, limit=ROUTER_CONTEXT_CHARS
        )

        try:
            with st.spinner("Analyzing task complexity…"):
                route = route_task(client, router_task)

            research = None
            if route.get("needs_research") and route.get("research_query", "").strip():
                if perplexity_key:
                    with st.spinner("Researching current information…"):
                        research = run_research(
                            perplexity_key, route["research_query"].strip()
                        )
                    if research:
                        notes.append("🔎 Grounded in live Perplexity research.")
                    else:
                        notes.append(
                            "🔎 Live research was unavailable — planned from "
                            "built-in knowledge."
                        )
                else:
                    notes.append(
                        "🔎 This task would benefit from live research — add a "
                        "Perplexity API key to enable it."
                    )

            # Non-simple tasks: draft 2-3 genuinely different candidate
            # workflows and keep the simplest one that fully delivers the
            # goal. Best-effort — if the candidate pass fails, fall back to
            # the single-workflow path below.
            draft = None
            selection_note = None
            if route["complexity"] == "complex":
                try:
                    with st.spinner("Designing candidate workflows…"):
                        payload = generate_candidate_workflows(
                            client, planner_task, selected_tools, research=research
                        )
                    draft, n_candidates, selection_note = pick_selected_candidate(
                        payload
                    )
                    notes.append(
                        f"⚖️ Compared {n_candidates} candidate workflows — "
                        "presenting the simplest that fully delivers the goal."
                    )
                except (
                    anthropic.APIError,
                    json.JSONDecodeError,
                    StopIteration,
                    KeyError,
                    IndexError,
                    TypeError,
                ):
                    draft = None

            if draft is None:
                with st.spinner("Designing your workflow…"):
                    draft, _ = generate_workflow(
                        client, planner_task, route["complexity"], selected_tools,
                        research=research,
                    )
        except anthropic.AuthenticationError:
            st.error("Invalid Anthropic API key.")
            st.stop()
        except anthropic.RateLimitError:
            st.error("The planner is busy right now — please try again in a moment.")
            st.stop()
        except anthropic.APIStatusError as e:
            st.error(f"The planner request failed ({e.status_code}). Please try again.")
            st.stop()
        except anthropic.APIConnectionError:
            st.error("Could not reach the planning service. Check your connection.")
            st.stop()
        except (json.JSONDecodeError, StopIteration, KeyError):
            st.error("The planner returned an unexpected response. Please try again.")
            st.stop()

        # The final review is a best-effort quality pass — if it can't complete
        # (rate limit, API error, malformed response), fall back to the
        # unreviewed draft rather than failing the whole run.
        try:
            with st.spinner("Running the final alignment review…"):
                workflow = review_workflow(
                    client, planner_task, selected_tools, draft,
                    research=research, selection_note=selection_note,
                )
        except (
            anthropic.APIError,
            json.JSONDecodeError,
            StopIteration,
            KeyError,
        ):
            workflow = draft
            notes.append(
                "⚠️ The final review step didn't complete — showing the "
                "unreviewed draft plan."
            )

        st.session_state["workflow"] = workflow
        st.session_state["workflow_task"] = task.strip()
        st.session_state["workflow_notes"] = notes
        st.session_state["workflow_attachments"] = attachment_names
        st.session_state["workflow_context_sig"] = context_sig

    if "workflow" in st.session_state:
        st.divider()
        for note in st.session_state.get("workflow_notes", []):
            st.caption(note)
        render_workflow(
            st.session_state.get("workflow_task", ""),
            st.session_state["workflow"],
            attachments=st.session_state.get("workflow_attachments") or None,
        )


if __name__ == "__main__":
    main()
